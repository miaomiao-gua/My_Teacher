import json
import os
import re
import sys
from pathlib import Path
from typing import Any, Dict, List
from urllib.parse import urlparse

import requests
from docx import Document
from pypdf import PdfReader

# python-pptx 安装在项目本地目录（D 盘，避免写入系统 site-packages 受限）
_PY_DEPS = Path(__file__).resolve().parent / "py_deps"
if _PY_DEPS.exists() and str(_PY_DEPS) not in sys.path:
    sys.path.insert(0, str(_PY_DEPS))
from pptx import Presentation  # noqa: E402


ALLOWED_SCHEMES = {"http", "https"}
ALLOWED_EXTENSIONS = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".txt", ".md", ".html"}


def sanitize_topic(topic: str) -> str:
    cleaned = re.sub(r"[^a-zA-Z0-9\u4e00-\u9fa5]+", "_", topic.strip())
    return cleaned.strip("_") or "lesson"


def ensure_lesson_dir(lesson_name: str, base_dir: str | Path | None = None) -> Path:
    base_path = Path(base_dir) if base_dir else Path(__file__).resolve().parent / "lessons"
    lesson_path = base_path / lesson_name
    lesson_path.mkdir(parents=True, exist_ok=True)
    return lesson_path


def safe_url_check(url: str) -> bool:
    parsed = urlparse(url)
    if parsed.scheme.lower() not in ALLOWED_SCHEMES:
        return False
    return True


def download_resource(resource: Dict[str, Any], lesson_dir: str | Path, file_prefix: str = "resource") -> Path:
    url = resource.get("url")
    if not url or not safe_url_check(url):
        raise ValueError(f"Invalid or unsafe URL: {url!r}")

    # 请求日志：记录向谁发请求包
    print(f"[DOWNLOAD] 发起下载请求: {url} | 标题: {resource.get('title', 'unnamed')} | 目标: {lesson_dir}", flush=True)

    lesson_path = Path(lesson_dir)
    lesson_path.mkdir(parents=True, exist_ok=True)

    try:
        response = requests.get(url, timeout=30)
        response.raise_for_status()
        print(f"[DOWNLOAD] 下载成功: {url} → HTTP {response.status_code}, {len(response.content)} bytes", flush=True)
    except Exception as exc:
        print(f"[DOWNLOAD] 下载失败: {url} → {type(exc).__name__}: {exc}", flush=True)
        raise

    content_type = response.headers.get("Content-Type", "")
    if "pdf" not in content_type and "doc" not in content_type and "text" not in content_type and "html" not in content_type:
        # Some documents may still be valid even if the server misses the content type.
        pass

    suffix = ".pdf"
    if "docx" in content_type.lower() or url.lower().endswith(".docx"):
        suffix = ".docx"
    elif "word" in content_type.lower() or url.lower().endswith(".doc"):
        suffix = ".doc"
    elif "pptx" in content_type.lower() or url.lower().endswith(".pptx"):
        suffix = ".pptx"
    elif url.lower().endswith(".txt"):
        suffix = ".txt"
    elif url.lower().endswith(".html") or url.lower().endswith(".htm"):
        suffix = ".html"

    target_path = lesson_path / f"{file_prefix}{suffix}"
    target_path.write_bytes(response.content)

    if suffix in {".pdf", ".docx", ".doc", ".pptx"}:
        markdown_path = target_path.with_suffix(".md")
        markdown_text = convert_document_to_markdown(target_path)
        markdown_path.write_text(markdown_text, encoding="utf-8")

    return target_path


def _read_pdf_text(file_path: str | Path, start_page: int | None = None, end_page: int | None = None) -> str:
    """提取 PDF 文本，支持页码范围（页码从 1 开始，含两端）。

    文字提取为空时回退 OCR（若本地可用）；OCR 也失败则返回空字符串。
    """
    try:
        reader = PdfReader(str(file_path))
        total = len(reader.pages)
        s = 1 if start_page is None else max(1, int(start_page))
        e = total if end_page is None else min(total, int(end_page))
        pages = []
        for i in range(s - 1, e):
            pages.append(reader.pages[i].extract_text() or "")
        text = "\n\n".join(pages)
        if text.strip():
            return text
        # 扫描版（提取不到文字）→ OCR 兜底
        return _pdf_ocr_pages(str(file_path), s, e)
    except Exception:
        return ""


# OCR 可用性缓存：避免每次请求都重复导入/检测（easyocr 导入较慢）
_OCR_CACHE = {"checked": False, "available": False}
# OCR 总开关（由 config.json 的 ocr_enabled 控制，app.py 启动/保存配置时同步）
_OCR_ENABLED = True


def set_ocr_enabled(enabled: bool) -> None:
    """设置 OCR 总开关（来自 config.json 的 ocr_enabled 字段）。"""
    global _OCR_ENABLED
    _OCR_ENABLED = bool(enabled)


def pdf_ocr_available() -> bool:
    """检测本地 OCR 是否可用（easyocr + 完整模型文件）。结果进程内缓存。"""
    if not _OCR_ENABLED:
        return False
    if _OCR_CACHE["checked"]:
        return _OCR_CACHE["available"]
    _OCR_CACHE["checked"] = True
    try:
        import os
        import zipfile
        import easyocr  # noqa: F401
        # 模型目录 ~/.EasyOCR/model 下应有完整模型（detection + recognition）。
        # zip 需校验完整性（损坏/半截下载的模型视为不可用）；pth 需有实际大小。
        model_dir = os.path.join(os.path.expanduser("~"), ".EasyOCR", "model")
        if os.path.isdir(model_dir):
            good = 0
            for fname in os.listdir(model_dir):
                fp = os.path.join(model_dir, fname)
                try:
                    if fname.endswith(".zip"):
                        with open(fp, "rb") as fh:
                            head = fh.read(4)
                        if head == b"PK\x03\x04" and zipfile.is_zipfile(fp) and os.path.getsize(fp) > 100_000:
                            good += 1
                    elif fname.endswith(".pth") and os.path.getsize(fp) > 500_000:
                        good += 1
                except Exception:
                    continue
            _OCR_CACHE["available"] = good >= 2  # 至少检测 + 识别模型各一
    except Exception:
        pass
    return _OCR_CACHE["available"]


def _pdf_ocr_pages(file_path: str | Path, start_page: int, end_page: int) -> str:
    """对扫描版 PDF 页码范围做 OCR：fitz 渲染页 → easyocr 识别。"""
    if not pdf_ocr_available():
        return ""
    try:
        import fitz  # PyMuPDF：把 PDF 页渲染成图片（无 poppler 依赖）
        import easyocr
        import numpy as np
        from PIL import Image
    except Exception:
        return ""
    try:
        reader = fitz.open(str(file_path))
        # easyocr 实例较重，懒加载单例
        if not hasattr(_pdf_ocr_pages, "_reader"):
            _pdf_ocr_pages._reader = easyocr.Reader(["ch_sim", "en"], gpu=False, verbose=False)
        ocr = _pdf_ocr_pages._reader
        parts = []
        for pno in range(max(0, start_page - 1), min(reader.page_count, end_page)):
            page = reader.load_page(pno)
            pix = page.get_pixmap(dpi=150)
            img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
            result = ocr.readtext(np.array(img), detail=0, paragraph=True)
            if result:
                parts.append(f"\n\n".join(result))
        reader.close()
        return "\n\n".join(parts)
    except Exception:
        return ""


def extract_pdf_toc(file_path: str | Path) -> List[Dict[str, Any]]:
    """提取 PDF 书签目录（目录）。返回 [{title, page, level}]，page 为 1 起始页码。

    无书签的 PDF 返回空列表（不强行猜测目录）。
    """
    try:
        reader = PdfReader(str(file_path))
        outline = reader.outline
        if not outline:
            return []
        items: List[Dict[str, Any]] = []

        def _walk(node, level):
            for item in node:
                if isinstance(item, list):
                    _walk(item, level + 1)
                    continue
                title = str(getattr(item, "title", "") or "").strip()
                if not title:
                    continue
                try:
                    page = reader.get_destination_page_number(item) + 1
                except Exception:
                    page = 0
                items.append({"title": title, "page": page, "level": level})

        _walk(outline, 0)
        return items
    except Exception:
        return []


def _read_docx_text(file_path: str | Path) -> str:
    try:
        doc = Document(str(file_path))
        return "\n".join(paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip())
    except Exception:
        return ""


def _read_pptx_text(file_path: str | Path) -> str:
    """读取 .pptx：按幻灯片顺序提取所有文本框内容（含表格）。"""
    try:
        prs = Presentation(str(file_path))
        slides: List[str] = []
        for idx, slide in enumerate(prs.slides, 1):
            parts: List[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = "\n".join(
                        para.text for para in shape.text_frame.paragraphs if para.text.strip()
                    )
                    if text.strip():
                        parts.append(text)
                if getattr(shape, "has_table", False) and shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        if any(cells):
                            parts.append(" | ".join(cells))
            if parts:
                slides.append(f"## 幻灯片 {idx}\n\n" + "\n".join(parts))
        return "\n\n".join(slides)
    except Exception:
        return ""


def convert_document_to_markdown(file_path: str | Path, start_page: int | None = None, end_page: int | None = None) -> str:
    path = Path(file_path)
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        text = _read_pdf_text(path, start_page, end_page)
    elif suffix in {".docx", ".doc"}:
        text = _read_docx_text(path)
    elif suffix == ".pptx":
        text = _read_pptx_text(path)
    elif suffix == ".ppt":
        text = ""
    elif suffix in {".txt", ".md"}:
        text = path.read_text(encoding="utf-8", errors="ignore")
    else:
        text = path.read_text(encoding="utf-8", errors="ignore")

    if not text.strip():
        if suffix == ".ppt":
            return "# 文档内容\n\n无法直接读取旧版 .ppt（二进制格式）。请用 PowerPoint 另存为 .pptx 后重新上传。"
        return "# 文档内容\n\n无法自动提取正文内容。"

    lines = [line.strip() for line in text.splitlines() if line.strip()]
    sections = []
    for line in lines:
        sections.append(line)
    return "\n\n".join(f"{item}" for item in sections)


def load_course_context(lesson_folder: str | None) -> str:
    """读取整门课程根目录下的 .md/.txt（兼容旧课程结构）。"""
    if not lesson_folder:
        return "未选择课程，使用通用教学模式。"

    lesson_dir = Path(__file__).resolve().parent / "lessons" / lesson_folder
    if not lesson_dir.exists():
        return "该课程目录不存在。"

    parts: List[str] = []
    for md_file in sorted(lesson_dir.glob("*.md")):
        try:
            parts.append(md_file.read_text(encoding="utf-8", errors="ignore"))
        except Exception:
            continue
    for txt_file in sorted(lesson_dir.glob("*.txt")):
        try:
            parts.append(txt_file.read_text(encoding="utf-8", errors="ignore"))
        except Exception:
            continue

    context = "\n\n---\n\n".join(parts)
    return context[:20000]


def unit_dir(lesson_folder: str, unit_index: int) -> Path:
    """返回某课对应的资料目录路径：lessons/<lesson>/section_N/。"""
    return Path(__file__).resolve().parent / "lessons" / lesson_folder / f"section_{unit_index + 1}"


def load_unit_context(lesson_folder: str | None, unit_index: int) -> str:
    """只读取指定 section 目录下的 .md/.txt 作为本课的资料正文。

    兼容旧结构：若 section_N 不存在，回退到 units/unit_NN/。
    """
    if not lesson_folder:
        return ""
    # 新结构：section_N
    udir = unit_dir(lesson_folder, unit_index)
    # 回退：旧结构 units/unit_NN
    if not udir.exists():
        udir_old = Path(__file__).resolve().parent / "lessons" / lesson_folder / "units" / f"unit_{unit_index + 1:02d}"
        if udir_old.exists():
            udir = udir_old
        else:
            return ""
    parts: List[str] = []
    for md_file in sorted(list(udir.glob("*.md")) + list(udir.glob("*.txt"))):
        try:
            parts.append(md_file.read_text(encoding="utf-8", errors="ignore"))
        except Exception:
            continue
    context = "\n\n---\n\n".join(parts)
    return context[:20000]


def save_metadata(lesson_dir: str | Path, payload: Dict[str, Any]) -> None:
    lesson_path = Path(lesson_dir)
    lesson_path.mkdir(parents=True, exist_ok=True)
    # 新结构存 config.json，兼容旧结构的 metadata.json
    config_path = lesson_path / "config.json"
    # 合并已有配置（保留 tools 等字段）
    existing = {}
    if config_path.exists():
        try:
            existing = json.loads(config_path.read_text(encoding="utf-8"))
        except Exception:
            existing = {}
    existing.update(payload)
    with open(config_path, "w", encoding="utf-8") as f:
        json.dump(existing, f, ensure_ascii=False, indent=2)
